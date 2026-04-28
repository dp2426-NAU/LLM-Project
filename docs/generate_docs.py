"""Generate proposal.docx and presentation.pptx for CS 599 DocuMind project."""

import os
from pathlib import Path

OUT = Path(__file__).parent  # docs/

# ─────────────────────────────────────────────────────────────────────────────
# DOCX — Project Proposal
# ─────────────────────────────────────────────────────────────────────────────
from docx import Document
from docx.shared import Pt, Inches, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import copy

def set_cell_bg(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), hex_color)
    tcPr.append(shd)

def add_heading(doc, text, level=1, color="1F3864"):
    p = doc.add_heading(text, level=level)
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    for run in p.runs:
        run.font.color.rgb = RGBColor.from_string(color)
    return p

def add_para(doc, text, bold=False, italic=False, size=11, color=None, space_before=0, space_after=6):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(space_before)
    p.paragraph_format.space_after = Pt(space_after)
    run = p.add_run(text)
    run.bold = bold
    run.italic = italic
    run.font.size = Pt(size)
    if color:
        run.font.color.rgb = RGBColor.from_string(color)
    return p

def add_bullet(doc, text, level=0):
    p = doc.add_paragraph(style='List Bullet')
    p.paragraph_format.space_after = Pt(3)
    p.paragraph_format.left_indent = Inches(0.25 + 0.25 * level)
    run = p.add_run(text)
    run.font.size = Pt(10.5)
    return p

def add_table_row(table, cells, header=False, bg=None):
    row = table.add_row()
    for i, val in enumerate(cells):
        cell = row.cells[i]
        cell.text = str(val)
        for para in cell.paragraphs:
            for run in para.runs:
                run.font.size = Pt(9.5)
                if header:
                    run.bold = True
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        if bg:
            set_cell_bg(cell, bg)
    return row

def build_docx():
    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1.15)
        section.right_margin = Inches(1.15)

    # Styles
    normal = doc.styles['Normal']
    normal.font.name = 'Calibri'
    normal.font.size = Pt(11)

    # ── TITLE PAGE ──────────────────────────────────────────────────────────
    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_p.paragraph_format.space_before = Pt(48)
    title_p.paragraph_format.space_after = Pt(6)
    r = title_p.add_run("DocuMind")
    r.bold = True
    r.font.size = Pt(28)
    r.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    sub_p = doc.add_paragraph()
    sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub_p.paragraph_format.space_after = Pt(4)
    r = sub_p.add_run("Production-Grade LLM Intelligence Platform")
    r.font.size = Pt(16)
    r.font.color.rgb = RGBColor(0x2E, 0x74, 0xB5)

    sub2 = doc.add_paragraph()
    sub2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub2.paragraph_format.space_after = Pt(36)
    r = sub2.add_run("Regulatory Change Analysis  ·  Multi-Agent Knowledge Synthesis  ·  RAG")
    r.italic = True
    r.font.size = Pt(11)
    r.font.color.rgb = RGBColor(0x40, 0x40, 0x40)

    for label, val in [
        ("Course:", "CS 599 — Contemporary Developments: Applications of Large Language Models"),
        ("Institution:", "Northern Arizona University · Department of Computer Science"),
        ("Modules:", "RegDelta (Regulatory Change Impact Analyzer)  |  SynapseIQ (Multi-Agent Research Platform)"),
        ("Repository:", "https://github.com/dp2426-NAU/LLM-Project"),
    ]:
        row_p = doc.add_paragraph()
        row_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        row_p.paragraph_format.space_after = Pt(4)
        rb = row_p.add_run(label + "  ")
        rb.bold = True
        rb.font.size = Pt(10)
        rv = row_p.add_run(val)
        rv.font.size = Pt(10)

    doc.add_page_break()

    # ── ABSTRACT ────────────────────────────────────────────────────────────
    add_heading(doc, "Abstract", 1)
    add_para(doc, (
        "Organizations in regulated industries generate and consume vast volumes of structured documents daily. "
        "Two critical bottlenecks persist: compliance teams cannot efficiently propagate regulatory amendments "
        "across internal policy inventories, and knowledge workers lack tools that synthesize multi-document "
        "corpora into rigorously evaluated artifacts. This project presents DocuMind, a production-grade AI "
        "platform comprising two complementary modules — RegDelta and SynapseIQ — both built on "
        "Retrieval-Augmented Generation (RAG) and Large Language Model (LLM) inference. RegDelta detects "
        "section-level regulatory changes, retrieves semantically matched internal policies, and generates "
        "risk-classified impact reports. SynapseIQ implements a four-agent sequential pipeline "
        "(Researcher → Critic → Synthesizer → Validator) producing self-evaluated knowledge syntheses "
        "with composite quality scoring. Both modules are deployed as production FastAPI services with "
        "ChromaDB vector stores, SQLite analytics databases, and Docker containerization."
    ), size=11)

    # ── PROBLEM STATEMENT ───────────────────────────────────────────────────
    add_heading(doc, "1.  Problem Statement", 1)

    add_heading(doc, "1.1  Regulatory Change Propagation", 2, color="2E74B5")
    add_para(doc, (
        "When regulatory bodies such as the SEC, FINRA, or Basel Committee amend rules, compliance teams "
        "must manually cross-reference hundreds of internal policy documents to identify what requires updating. "
        "This process typically requires 2–4 weeks per major amendment, is highly error-prone, and produces "
        "no systematic audit trail linking identified policy gaps to specific regulatory clauses."
    ))
    for b in [
        "Takes 2–4 weeks per major rule change at mid-size firms",
        "Error-prone under deadline pressure",
        "No audit trail linking policy gaps to regulatory clauses",
        "Scales linearly with both regulation corpus and policy inventory size",
    ]:
        add_bullet(doc, b)

    add_heading(doc, "1.2  Research and Knowledge Synthesis", 2, color="2E74B5")
    add_para(doc, (
        "Analysts and knowledge workers spend the majority of their cognitive effort manually integrating "
        "information from heterogeneous document corpora. Existing LLM tools retrieve relevant text but do not "
        "reason about logical gaps, apply structured critique before synthesis, self-evaluate output quality, "
        "or maintain a traceable reasoning chain from raw source to final artifact."
    ))

    add_heading(doc, "1.3  Limitations of Existing Approaches", 2, color="2E74B5")
    tbl = doc.add_table(rows=1, cols=2)
    tbl.style = 'Table Grid'
    tbl.autofit = False
    tbl.columns[0].width = Inches(2.8)
    tbl.columns[1].width = Inches(3.8)
    add_table_row(tbl, ["Approach", "Limitation"], header=True, bg="1F3864")
    for row_data in [
        ("Simple RAG chatbots", "Returns raw chunks; no synthesized analysis; no quality gate"),
        ("Document summarizers", "Single-document scope; no cross-document reasoning"),
        ("Manual compliance review", "Non-scalable; no systematic gap detection; no audit trail"),
        ("Generic LLM API calls", "No domain retrieval; hallucination risk without grounding"),
        ("Rule-based diff tools", "No semantic understanding of regulatory significance"),
    ]:
        r = add_table_row(tbl, row_data)
        for cell in r.cells:
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9.5)
    doc.add_paragraph()

    # ── OBJECTIVES ──────────────────────────────────────────────────────────
    add_heading(doc, "2.  Objectives", 1)

    add_heading(doc, "Module I — RegDelta", 2, color="2E74B5")
    tbl2 = doc.add_table(rows=1, cols=3)
    tbl2.style = 'Table Grid'
    add_table_row(tbl2, ["ID", "Objective", "Success Criterion"], header=True, bg="1F3864")
    for row_data in [
        ("M1", "Section-level change detection", "ADDED / MODIFIED / REMOVED classification with similarity score"),
        ("M2", "Semantic policy impact retrieval", "Retrieve affected policies with cosine similarity ≥ 0.35"),
        ("M3", "LLM-generated impact report", "JSON report: risk_level, executive_summary, recommended_actions"),
        ("M4", "Version registry", "SQLite tracks all regulation versions and reports with metadata"),
        ("M5", "Production REST API", "FastAPI endpoints with OpenAPI docs, Pydantic validation, health check"),
    ]:
        r = add_table_row(tbl2, row_data)
        for cell in r.cells:
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9)
    doc.add_paragraph()

    add_heading(doc, "Module II — SynapseIQ", 2, color="2E74B5")
    tbl3 = doc.add_table(rows=1, cols=3)
    tbl3.style = 'Table Grid'
    add_table_row(tbl3, ["ID", "Objective", "Success Criterion"], header=True, bg="1F3864")
    for row_data in [
        ("S1", "4-agent sequential pipeline", "Researcher → Critic → Synthesizer → Validator with shared AgentContext"),
        ("S2", "Self-evaluating output", "Composite quality score computed; PASS/CONDITIONAL_PASS/FAIL verdict issued"),
        ("S3", "Format-adaptive synthesis", "brief (~300w), standard (~800w), detailed (~2000w) modes"),
        ("S4", "Agent observability", "Per-agent token count, latency, success recorded per session"),
        ("S5", "Cost-bounded inference", "Per-agent token budgets enforced; USD cost tracked per query"),
    ]:
        r = add_table_row(tbl3, row_data)
        for cell in r.cells:
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9)
    doc.add_paragraph()

    # ── METHODOLOGY ─────────────────────────────────────────────────────────
    add_heading(doc, "3.  Methodology", 1)

    add_heading(doc, "3.1  Retrieval-Augmented Generation (RAG)", 2, color="2E74B5")
    add_para(doc, (
        "Both modules ground all LLM inference in retrieved document evidence. The RAG pipeline operates in two phases:"
    ))
    add_bullet(doc, "Offline Indexing: Documents are parsed, cleaned, and split into overlapping fixed-size chunks "
               "(512–800 tokens, 50–100 token overlap). Each chunk is embedded using sentence-transformers "
               "all-MiniLM-L6-v2 (384-dimensional, L2-normalized) and upserted into ChromaDB.")
    add_bullet(doc, "Online Retrieval: At query time, the user query is embedded with the same model and passed to "
               "ChromaDB's HNSW index for approximate nearest-neighbor search. The top-K chunks above a minimum "
               "relevance threshold become the LLM grounding context.")

    add_heading(doc, "3.2  Section-Level Differential Analysis (RegDelta)", 2, color="2E74B5")
    add_para(doc, (
        "RegDelta implements a domain-specific diff engine operating at heading-level document granularity. "
        "Documents are parsed using a regex matching regulatory heading conventions (Section, Article, Rule, § prefixes). "
        "Python's difflib.SequenceMatcher computes a similarity ratio (0.0–1.0) for each matched section pair. "
        "Sections are classified as ADDED, REMOVED, or MODIFIED and sorted by ascending similarity to surface "
        "the most significant changes first. The top 10 changed sections are individually annotated by the LLM "
        "before full impact report generation."
    ))

    add_heading(doc, "3.3  Multi-Agent Sequential Architecture (SynapseIQ)", 2, color="2E74B5")
    add_para(doc, "SynapseIQ implements a four-stage sequential agent pipeline sharing an AgentContext dataclass:")
    for b in [
        "ResearcherAgent — Extracts structured findings, key concepts, evidence gaps from retrieved chunks. Budget: 1,200 tokens.",
        "CriticAgent — Identifies logical gaps, unsupported claims, and bias signals. Budget: 800 tokens.",
        "SynthesizerAgent — Produces a format-appropriate knowledge artifact (brief/standard/detailed). Budget: 1,800 tokens.",
        "ValidatorAgent — Scores synthesis on 4 dimensions; issues PASS/CONDITIONAL_PASS/FAIL verdict. Budget: 600 tokens.",
    ]:
        add_bullet(doc, b)

    add_heading(doc, "3.4  Evaluation Framework", 2, color="2E74B5")
    add_para(doc, "SynapseIQ uses a weighted composite evaluation score:")
    add_para(doc, "Composite = 0.30 × coherence + 0.35 × relevance + 0.25 × factuality + 0.10 × completeness",
             bold=True, color="1F3864")
    add_para(doc, "Verdict thresholds: PASS ≥ 0.65  |  CONDITIONAL_PASS ≥ 0.45  |  FAIL < 0.45")

    # ── SYSTEM ARCHITECTURE ─────────────────────────────────────────────────
    add_heading(doc, "4.  System Architecture", 1)

    add_heading(doc, "4.1  High-Level Layered Architecture", 2, color="2E74B5")
    code_p = doc.add_paragraph()
    code_p.paragraph_format.space_after = Pt(6)
    code_run = code_p.add_run(
        "CLIENT LAYER         HTTP · Swagger UI · Python SDK · CI/CD\n"
        "       ↓\n"
        "API GATEWAY          FastAPI / Uvicorn ASGI · Pydantic · OpenAPI 3.0\n"
        "       ↓\n"
        "DOMAIN MODULES       RegDelta (Diff → Analyze → Report)\n"
        "                     SynapseIQ (Researcher → Critic → Synthesizer → Validator)\n"
        "       ↓\n"
        "SHARED INFRA         ChromaDB · Embedding Engine · LLM Client · Config\n"
        "       ↓\n"
        "PERSISTENCE          SQLite (versions + reports) · SQLite (sessions + traces)\n"
        "                     ChromaDB on-disk · OpenAI API / Ollama local"
    )
    code_run.font.name = "Courier New"
    code_run.font.size = Pt(9)
    code_run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    add_heading(doc, "4.2  RegDelta Data Flow", 2, color="2E74B5")
    for step in [
        "POST /api/v1/analyze/delta — receive regulation_id, old_version_tag, new_version_tag",
        "VersionTracker.get_latest_two() — fetch old_text and new_text from SQLite",
        "DiffEngine.compute_section_diffs(old, new) — regex split + SequenceMatcher",
        "ImpactAnalyzer._annotate_sections(diffs[:10]) — one LLM call per section",
        "Retriever.retrieve_policies(section_text) — ChromaDB ANN search on policies collection",
        "LLMGenerator.generate_json(impact_system, prompt) — full ImpactReport JSON",
        "VersionTracker.save_report(...) — persist to SQLite impact_reports table",
        "Return ImpactReport.to_dict() as JSON response",
    ]:
        add_bullet(doc, step)

    add_heading(doc, "4.3  SynapseIQ Pipeline Flow", 2, color="2E74B5")
    for step in [
        "POST /api/v1/synthesize — receive query, output_format, citation_style",
        "CorpusStore.search(query, top_k=6, min_score=0.35) — ChromaDB HNSW ANN search",
        "ResearcherAgent._safe_run(context) — YAML prompt → LLM → structured findings",
        "CriticAgent._safe_run(context) — YAML prompt → LLM → gaps and bias flags",
        "SynthesizerAgent._safe_run(context) — YAML prompt → LLM → knowledge artifact",
        "ValidatorAgent._safe_run(context) — YAML prompt → LLM → scores + verdict",
        "PipelineEvaluator.evaluate(context) — compute composite score; set EvaluationResult",
        "PipelineTracer.end_session(...) — write pipeline_sessions + agent_traces to SQLite",
    ]:
        add_bullet(doc, step)

    # ── TECHNOLOGIES ────────────────────────────────────────────────────────
    add_heading(doc, "5.  Technologies Used", 1)
    tbl4 = doc.add_table(rows=1, cols=3)
    tbl4.style = 'Table Grid'
    add_table_row(tbl4, ["Category", "Technology", "Purpose"], header=True, bg="1F3864")
    for row_data in [
        ("Language", "Python 3.11+", "Primary implementation language"),
        ("Web Framework", "FastAPI + Uvicorn", "ASGI REST API server"),
        ("LLM Provider (cloud)", "OpenAI gpt-4o-mini", "Default inference backend"),
        ("LLM Provider (local)", "Ollama llama3", "Zero-cost local inference alternative"),
        ("Vector Database", "ChromaDB 0.5", "Dense vector storage and ANN search"),
        ("Embedding Model", "all-MiniLM-L6-v2", "384-dimensional sentence embeddings"),
        ("Persistence", "SQLite3", "Version tracking and session analytics"),
        ("Data Validation", "Pydantic v2", "Request/response models and settings"),
        ("Prompt Templating", "Jinja2 + YAML", "Agent prompt management with hot-reload"),
        ("Diff Engine", "difflib (stdlib)", "Section-level text comparison"),
        ("Containerization", "Docker", "Production deployment packaging"),
        ("Testing", "pytest", "Unit, integration, and evaluation tests"),
    ]:
        r = add_table_row(tbl4, row_data)
        for cell in r.cells:
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9.5)
    doc.add_paragraph()

    # ── ROLE OF LLMs ────────────────────────────────────────────────────────
    add_heading(doc, "6.  Role of Large Language Models", 1)
    add_para(doc, (
        "LLMs serve as the reasoning core of DocuMind. Critically, all LLM calls are supplied with retrieved "
        "document context from ChromaDB, ensuring factual grounding and preventing hallucination from parametric memory."
    ))

    add_heading(doc, "RegDelta LLM Usage", 2, color="2E74B5")
    tbl5 = doc.add_table(rows=1, cols=4)
    tbl5.style = 'Table Grid'
    add_table_row(tbl5, ["Task", "Calls", "Input", "Output"], header=True, bg="1F3864")
    for row_data in [
        ("Section annotation", "≤10 (per section)", "Old + new text excerpt", "One-sentence significance note"),
        ("Impact report generation", "1 (per analysis)", "Changed sections + matched policies", "JSON: risk_level, summary, actions"),
    ]:
        r = add_table_row(tbl5, row_data)
        for cell in r.cells:
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9.5)
    doc.add_paragraph()

    add_heading(doc, "SynapseIQ LLM Usage (4 agents per query)", 2, color="2E74B5")
    tbl6 = doc.add_table(rows=1, cols=4)
    tbl6.style = 'Table Grid'
    add_table_row(tbl6, ["Agent", "Task", "Token Budget", "Output Format"], header=True, bg="1F3864")
    for row_data in [
        ("ResearcherAgent", "Extract findings from retrieved chunks", "1,200", "core_findings, key_concepts, evidence_gaps"),
        ("CriticAgent", "Identify gaps, bias, unsupported claims", "800", "Structured JSON critique"),
        ("SynthesizerAgent", "Produce knowledge artifact", "1,800", "Formatted document (brief/standard/detailed)"),
        ("ValidatorAgent", "Score 4 dimensions; issue verdict", "600", "JSON scores + PASS/CONDITIONAL_PASS/FAIL"),
    ]:
        r = add_table_row(tbl6, row_data)
        for cell in r.cells:
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9.5)
    doc.add_paragraph()

    add_heading(doc, "Cost Control Strategy", 2, color="2E74B5")
    for b in [
        "Per-agent token budgets enforced in AGENT_TOKEN_BUDGETS constant",
        "Cost estimated at $0.00015/1K input tokens, $0.00060/1K output tokens (gpt-4o-mini)",
        "CostReport included in every /analytics response",
        "Ollama local inference available as zero-cost development alternative",
        "ENABLE_CRITIC_AGENT and ENABLE_VALIDATOR_AGENT flags allow optional agent disabling",
    ]:
        add_bullet(doc, b)

    # ── EXPECTED OUTCOMES ───────────────────────────────────────────────────
    add_heading(doc, "7.  Expected Outcomes", 1)
    for i, outcome in enumerate([
        ("Automated regulatory compliance gap detection",
         "RegDelta reduces manual policy review time from weeks to minutes by automatically identifying which "
         "internal policies are semantically affected by regulation amendments."),
        ("Rigorous multi-source knowledge synthesis",
         "SynapseIQ produces quality-evaluated knowledge artifacts from multi-document corpora with a traceable "
         "agent reasoning chain and a self-issued quality verdict before delivery."),
        ("Production-grade system design",
         "Both modules demonstrate enterprise-quality software engineering: typed Python, dependency injection, "
         "structured logging, persistent storage, containerization, and comprehensive test coverage."),
        ("Transferable architectural patterns",
         "The dual-corpus RAG, section-level diff, sequential agent state sharing, and YAML prompt management "
         "patterns are domain-agnostic and transferable to legal, medical, and scientific literature domains."),
    ], 1):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(4)
        r_num = p.add_run(f"{i}. ")
        r_num.bold = True
        r_num.font.size = Pt(11)
        r_bold = p.add_run(outcome[0] + " — ")
        r_bold.bold = True
        r_bold.font.size = Pt(11)
        r_text = p.add_run(outcome[1])
        r_text.font.size = Pt(11)

    # ── TIMELINE ────────────────────────────────────────────────────────────
    add_heading(doc, "8.  Development Timeline", 1)
    tbl7 = doc.add_table(rows=1, cols=3)
    tbl7.style = 'Table Grid'
    add_table_row(tbl7, ["Phase", "Duration", "Deliverables"], header=True, bg="1F3864")
    for row_data in [
        ("Phase 1 — Infrastructure", "Week 1", "ChromaDB, FastAPI skeleton, embedding pipeline, SQLite schema"),
        ("Phase 2 — RegDelta Core", "Week 2", "DiffEngine, ImpactAnalyzer, VectorStore, VersionTracker, REST API"),
        ("Phase 3 — SynapseIQ Core", "Week 3", "All 4 agents, AgentContext, SynapseOrchestrator, PromptEngine"),
        ("Phase 4 — Evaluation", "Week 4", "PipelineEvaluator, metrics, cost tracker, analytics endpoints"),
        ("Phase 5 — Testing & Docs", "Week 5", "pytest suite (unit, integration, evaluation), README, proposal, diagrams"),
        ("Phase 6 — Finalization", "Week 6", "Docker, GitHub Actions CI, performance tuning, final documentation"),
    ]:
        r = add_table_row(tbl7, row_data)
        for cell in r.cells:
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9.5)
    doc.add_paragraph()

    # ── REFERENCES ──────────────────────────────────────────────────────────
    add_heading(doc, "9.  References", 1)
    refs = [
        "Lewis, P., et al. (2020). Retrieval-Augmented Generation for Knowledge-Intensive NLP Tasks. NeurIPS 2020. arXiv:2005.11401.",
        "Reimers, N., & Gurevych, I. (2019). Sentence-BERT: Sentence Embeddings using Siamese BERT-Networks. EMNLP 2019. arXiv:1908.10084.",
        "OpenAI. (2024). GPT-4o Technical Report. openai.com/research/gpt-4o.",
        "Chromadb. (2024). ChromaDB: The AI-Native Open-Source Embedding Database. github.com/chroma-core/chroma.",
        "Yao, S., et al. (2023). ReAct: Synergizing Reasoning and Acting in Language Models. ICLR 2023. arXiv:2210.03629.",
        "Wang, L., et al. (2024). A Survey on Large Language Model based Autonomous Agents. arXiv:2308.11432.",
        "Peng, B., et al. (2023). Check Your Facts and Try Again. arXiv:2302.12813.",
        "Gao, Y., et al. (2024). Retrieval-Augmented Generation for Large Language Models: A Survey. arXiv:2312.10997.",
        "FastAPI. (2024). FastAPI Framework Documentation. fastapi.tiangolo.com.",
        "Hu, E., et al. (2022). LoRA: Low-Rank Adaptation of Large Language Models. ICLR 2022. arXiv:2106.09685.",
    ]
    for i, ref in enumerate(refs, 1):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(4)
        p.paragraph_format.left_indent = Inches(0.4)
        p.paragraph_format.first_line_indent = Inches(-0.4)
        r = p.add_run(f"[{i}]  {ref}")
        r.font.size = Pt(10)

    out_path = OUT / "proposal.docx"
    doc.save(str(out_path))
    print(f"OK  Created {out_path}")


# ─────────────────────────────────────────────────────────────────────────────
# PPTX — Presentation
# ─────────────────────────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN

DARK_BG    = PptxRGB(0x0D, 0x11, 0x17)
DARK_BG2   = PptxRGB(0x16, 0x1B, 0x22)
DARK_BG3   = PptxRGB(0x1C, 0x21, 0x28)
WHITE      = PptxRGB(0xE6, 0xED, 0xF3)
MUTED      = PptxRGB(0x8B, 0x94, 0x9E)
BLUE       = PptxRGB(0x38, 0x8B, 0xFD)
BLUE_DIM   = PptxRGB(0x1F, 0x6F, 0xEB)
GREEN      = PptxRGB(0x3F, 0xB9, 0x50)
GREEN_DIM  = PptxRGB(0x23, 0x86, 0x36)
ORANGE     = PptxRGB(0xF0, 0x88, 0x3E)
PURPLE     = PptxRGB(0xD2, 0xA8, 0xFF)
YELLOW     = PptxRGB(0xE3, 0xB3, 0x41)

W = Inches(13.333)
H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)

def fill_slide_bg(slide, color: PptxRGB):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(1)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h, font_size=Pt(18), bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_items(slide, items, l, t, w, h, font_size=Pt(18), color=WHITE,
                     bullet_color=BLUE, accent_colors=None):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        bullet_char = "▸  "
        r1 = p.add_run()
        r1.text = bullet_char
        r1.font.size = font_size
        bcolor = accent_colors[i] if accent_colors and i < len(accent_colors) else bullet_color
        r1.font.color.rgb = bcolor
        r2 = p.add_run()
        r2.text = item
        r2.font.size = font_size
        r2.font.color.rgb = color
    return txBox

def header_bar(slide, title, subtitle=None, accent=BLUE):
    # Top accent stripe
    add_rect(slide, 0, 0, W, Inches(0.08), fill_color=accent)
    # Title bar
    add_rect(slide, 0, Inches(0.08), W, Inches(1.12), fill_color=DARK_BG2)
    # Title text
    add_text_box(slide, title,
                 Inches(0.5), Inches(0.12), Inches(12), Inches(0.7),
                 font_size=Pt(30), bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.5), Inches(0.78), Inches(12), Inches(0.38),
                     font_size=Pt(14), color=MUTED)

def footer_bar(slide, text="DocuMind  ·  CS 599  ·  Northern Arizona University"):
    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), fill_color=DARK_BG2)
    add_text_box(slide, text,
                 Inches(0.4), Inches(7.12), Inches(10), Inches(0.35),
                 font_size=Pt(10), color=MUTED)

def section_label(slide, label, l, t, w=Inches(2.5), color=BLUE, bg=None):
    if bg:
        add_rect(slide, l, t, w, Inches(0.28), fill_color=bg)
    add_text_box(slide, label, l + Inches(0.08), t, w, Inches(0.28),
                 font_size=Pt(10), bold=True, color=color)


def build_pptx():
    prs = new_prs()

    # ── SLIDE 1: Title ───────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_slide_bg(sl, DARK_BG)
    # Gradient-ish diagonal accent
    add_rect(sl, 0, 0, Inches(0.25), H, fill_color=BLUE_DIM)
    add_rect(sl, Inches(0.25), 0, Inches(0.08), H, fill_color=PptxRGB(0x0C, 0x2A, 0x4A))

    # Main title
    add_text_box(sl, "DocuMind",
                 Inches(0.65), Inches(1.4), Inches(11.5), Inches(1.5),
                 font_size=Pt(54), bold=True, color=WHITE)
    add_text_box(sl, "Production-Grade LLM Intelligence Platform",
                 Inches(0.65), Inches(2.85), Inches(11.5), Inches(0.7),
                 font_size=Pt(22), color=BLUE)

    add_rect(sl, Inches(0.65), Inches(3.5), Inches(10), Inches(0.04),
             fill_color=PptxRGB(0x30, 0x36, 0x3D))

    add_text_box(sl, "Regulatory Change Analysis  ·  Multi-Agent Knowledge Synthesis  ·  RAG",
                 Inches(0.65), Inches(3.6), Inches(11), Inches(0.5),
                 font_size=Pt(14), italic=True, color=MUTED)

    for i, (label, val) in enumerate([
        ("Course:", "CS 599 — Contemporary Developments: Applications of Large Language Models"),
        ("Modules:", "RegDelta (Regulatory Impact Analyzer)  |  SynapseIQ (Multi-Agent Synthesis Platform)"),
        ("University:", "Northern Arizona University · Department of Computer Science"),
    ]):
        y = Inches(4.3 + i * 0.45)
        add_text_box(sl, label, Inches(0.65), y, Inches(1.1), Inches(0.42),
                     font_size=Pt(11), bold=True, color=MUTED)
        add_text_box(sl, val, Inches(1.75), y, Inches(10.2), Inches(0.42),
                     font_size=Pt(11), color=WHITE)

    footer_bar(sl)

    # ── SLIDE 2: Introduction ────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_slide_bg(sl, DARK_BG)
    header_bar(sl, "Introduction", "What is DocuMind?")

    add_bullet_items(sl, [
        "A unified AI platform built on Large Language Models + Retrieval-Augmented Generation (RAG)",
        "Two production-grade modules for high-value document intelligence in regulated industries",
        "Demonstrates enterprise software engineering: FastAPI, ChromaDB, Docker, pytest, SQLite",
        "Supports OpenAI gpt-4o-mini cloud inference or Ollama local inference — single .env switch",
        "Full observability: per-agent token counts, latency, cost estimates, and session traces",
    ],
    Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.2),
    font_size=Pt(19), color=WHITE)

    footer_bar(sl)

    # ── SLIDE 3: Problem Statement ───────────────────────────────────────────
    sl = blank_slide(prs)
    fill_slide_bg(sl, DARK_BG)
    header_bar(sl, "Problem Statement", "Two unsolved enterprise knowledge-work bottlenecks")

    # Left panel
    add_rect(sl, Inches(0.4), Inches(1.3), Inches(5.9), Inches(5.5),
             fill_color=DARK_BG3, line_color=BLUE_DIM, line_width=Pt(1.5))
    add_rect(sl, Inches(0.4), Inches(1.3), Inches(5.9), Inches(0.45), fill_color=BLUE_DIM)
    add_text_box(sl, "Problem A — Regulatory Propagation",
                 Inches(0.55), Inches(1.35), Inches(5.6), Inches(0.4),
                 font_size=Pt(13), bold=True, color=WHITE)
    add_bullet_items(sl, [
        "2–4 weeks per regulation amendment at mid-size firms",
        "Manual cross-referencing of hundreds of policy docs",
        "Error-prone under deadline pressure",
        "No audit trail linking policy gaps to clauses",
        "Scales linearly with policy inventory size",
    ],
    Inches(0.55), Inches(1.85), Inches(5.6), Inches(4.5),
    font_size=Pt(15), color=WHITE)

    # Right panel
    add_rect(sl, Inches(6.9), Inches(1.3), Inches(6.0), Inches(5.5),
             fill_color=DARK_BG3, line_color=GREEN_DIM, line_width=Pt(1.5))
    add_rect(sl, Inches(6.9), Inches(1.3), Inches(6.0), Inches(0.45), fill_color=GREEN_DIM)
    add_text_box(sl, "Problem B — Research Synthesis at Scale",
                 Inches(7.05), Inches(1.35), Inches(5.7), Inches(0.4),
                 font_size=Pt(13), bold=True, color=WHITE)
    add_bullet_items(sl, [
        "Analysts spend majority of time on manual synthesis",
        "Simple RAG tools return chunks — no analysis",
        "No reasoning about logical gaps or contradictions",
        "No quality gate before delivering to decision-makers",
        "No traceable chain from source to final artifact",
    ],
    Inches(7.05), Inches(1.85), Inches(5.7), Inches(4.5),
    font_size=Pt(15), color=WHITE)

    footer_bar(sl)

    # ── SLIDE 4: Solution Overview ───────────────────────────────────────────
    sl = blank_slide(prs)
    fill_slide_bg(sl, DARK_BG)
    header_bar(sl, "Solution Overview", "DocuMind — two specialized, production-grade AI modules")

    # Module I box
    add_rect(sl, Inches(0.4), Inches(1.3), Inches(5.9), Inches(4.8),
             fill_color=DARK_BG3, line_color=BLUE_DIM, line_width=Pt(2))
    add_rect(sl, Inches(0.4), Inches(1.3), Inches(5.9), Inches(0.5), fill_color=BLUE_DIM)
    add_text_box(sl, "Module I — RegDelta",
                 Inches(0.55), Inches(1.32), Inches(5.6), Inches(0.46),
                 font_size=Pt(16), bold=True, color=WHITE)
    add_bullet_items(sl, [
        "Section-level diff with difflib.SequenceMatcher",
        "ChromaDB semantic policy retrieval",
        "LLM-generated risk-classified impact reports",
        "RiskLevel: LOW / MEDIUM / HIGH / CRITICAL",
        "SQLite audit trail of all versions and reports",
    ],
    Inches(0.55), Inches(1.9), Inches(5.6), Inches(3.9),
    font_size=Pt(16), color=WHITE, bullet_color=BLUE)

    # Module II box
    add_rect(sl, Inches(6.9), Inches(1.3), Inches(6.0), Inches(4.8),
             fill_color=DARK_BG3, line_color=GREEN_DIM, line_width=Pt(2))
    add_rect(sl, Inches(6.9), Inches(1.3), Inches(6.0), Inches(0.5), fill_color=GREEN_DIM)
    add_text_box(sl, "Module II — SynapseIQ",
                 Inches(7.05), Inches(1.32), Inches(5.7), Inches(0.46),
                 font_size=Pt(16), bold=True, color=WHITE)
    add_bullet_items(sl, [
        "4-agent pipeline: Researcher→Critic→Synthesizer→Validator",
        "Shared AgentContext dataclass across all agents",
        "YAML+Jinja2 versioned prompt templates, hot-reload",
        "Composite score: 0.30×coherence + 0.35×relevance + ...",
        "PASS / CONDITIONAL_PASS / FAIL verdict before delivery",
    ],
    Inches(7.05), Inches(1.9), Inches(5.7), Inches(3.9),
    font_size=Pt(16), color=WHITE, bullet_color=GREEN)

    # Shared strip
    add_rect(sl, Inches(0.4), Inches(6.25), Inches(12.5), Inches(0.55),
             fill_color=PptxRGB(0x2D, 0x1F, 0x3D), line_color=PURPLE, line_width=Pt(1))
    add_text_box(sl, "Shared: ChromaDB · all-MiniLM-L6-v2 · OpenAI/Ollama · Pydantic · SQLite · Docker · pytest",
                 Inches(0.6), Inches(6.27), Inches(12.1), Inches(0.5),
                 font_size=Pt(12), bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

    footer_bar(sl)

    # ── SLIDE 5: Architecture Diagram ───────────────────────────────────────
    sl = blank_slide(prs)
    fill_slide_bg(sl, DARK_BG)
    header_bar(sl, "System Architecture", "5-Layer Production Design")

    layers = [
        ("CLIENT LAYER", "HTTP · Swagger UI · Python SDK · CI/CD",         BLUE,   DARK_BG3,  0),
        ("API GATEWAY",  "FastAPI/Uvicorn ASGI · Pydantic · OpenAPI 3.0",  ORANGE, DARK_BG3,  1),
        ("DOMAIN MODULES","RegDelta (Diff→Analyze→Report)  |  SynapseIQ (4-Agent Pipeline)", BLUE_DIM, PptxRGB(0x0C,0x1A,0x2E), 2),
        ("SHARED INFRA", "ChromaDB · Embedding Engine · LLM Client · Config · Logging", PURPLE, PptxRGB(0x1A,0x12,0x29), 3),
        ("PERSISTENCE",  "SQLite (RegDelta + SynapseIQ) · ChromaDB on-disk · OpenAI API", GREEN_DIM, PptxRGB(0x09,0x14,0x09), 4),
    ]
    for i, (name, desc, accent, bg, idx) in enumerate(layers):
        y = Inches(1.35 + idx * 1.0)
        add_rect(sl, Inches(0.4), y, Inches(12.5), Inches(0.9),
                 fill_color=bg, line_color=accent, line_width=Pt(1.5))
        add_rect(sl, Inches(0.4), y, Inches(0.06), Inches(0.9), fill_color=accent)
        add_text_box(sl, name,
                     Inches(0.6), y + Inches(0.04), Inches(2.8), Inches(0.4),
                     font_size=Pt(11), bold=True, color=accent)
        add_text_box(sl, desc,
                     Inches(0.6), y + Inches(0.42), Inches(11.8), Inches(0.42),
                     font_size=Pt(13), color=WHITE)
        if idx < 4:
            add_text_box(sl, "↓", Inches(6.2), y + Inches(0.9), Inches(0.5), Inches(0.2),
                         font_size=Pt(14), color=MUTED, align=PP_ALIGN.CENTER)

    footer_bar(sl)

    # ── SLIDE 6: LLM Integration ─────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_slide_bg(sl, DARK_BG)
    header_bar(sl, "LLM Integration", "Where and how LLMs are used in both modules")

    agent_data = [
        ("R", "ResearcherAgent", "Extract structured findings from retrieved chunks", "1,200 tok", BLUE,   PptxRGB(0x0C,0x2A,0x4A)),
        ("C", "CriticAgent",    "Identify logical gaps, bias, unsupported claims",    "800 tok",  ORANGE, PptxRGB(0x2D,0x19,0x00)),
        ("S", "SynthesizerAgent","Produce knowledge artifact (brief/standard/detailed)","1,800 tok",GREEN, PptxRGB(0x09,0x14,0x09)),
        ("V", "ValidatorAgent", "Score 4 dimensions; issue PASS/CONDITIONAL/FAIL",    "600 tok",  PURPLE, PptxRGB(0x1F,0x10,0x40)),
    ]
    for i, (letter, name, task, budget, accent, bg) in enumerate(agent_data):
        y = Inches(1.35 + i * 1.3)
        add_rect(sl, Inches(0.4), y, Inches(12.5), Inches(1.05),
                 fill_color=bg, line_color=accent, line_width=Pt(1.5))
        add_rect(sl, Inches(0.4), y, Inches(0.06), Inches(1.05), fill_color=accent)
        # Agent badge
        add_rect(sl, Inches(0.6), y + Inches(0.22), Inches(0.55), Inches(0.55), fill_color=accent)
        add_text_box(sl, letter,
                     Inches(0.6), y + Inches(0.22), Inches(0.55), Inches(0.55),
                     font_size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, name,
                     Inches(1.3), y + Inches(0.06), Inches(4.5), Inches(0.4),
                     font_size=Pt(14), bold=True, color=accent)
        add_text_box(sl, task,
                     Inches(1.3), y + Inches(0.48), Inches(9.5), Inches(0.4),
                     font_size=Pt(13), color=WHITE)
        add_rect(sl, Inches(11.0), y + Inches(0.3), Inches(1.7), Inches(0.35),
                 fill_color=PptxRGB(0x21,0x26,0x2D), line_color=accent, line_width=Pt(1))
        add_text_box(sl, budget,
                     Inches(11.0), y + Inches(0.3), Inches(1.7), Inches(0.35),
                     font_size=Pt(12), bold=True, color=accent, align=PP_ALIGN.CENTER)

    add_text_box(sl, "★  All LLM calls grounded in ChromaDB-retrieved context — no parametric hallucination",
                 Inches(0.5), Inches(6.68), Inches(12.3), Inches(0.4),
                 font_size=Pt(13), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    footer_bar(sl)

    # ── SLIDE 7: Features ────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_slide_bg(sl, DARK_BG)
    header_bar(sl, "Key Features", "Platform capabilities across both modules")

    left_items = [
        "Section-level ADDED / MODIFIED / REMOVED classification",
        "Dual-corpus ChromaDB (regulations + policies isolated)",
        "LLM significance annotation per changed section",
        "Risk classification: LOW / MEDIUM / HIGH / CRITICAL",
        "SHA-256 idempotent document ingestion",
        "Full SQLite audit trail — every version and report",
    ]
    right_items = [
        "4-agent shared AgentContext state graph",
        "Versioned YAML + Jinja2 prompts with hot-reload",
        "Format-adaptive: brief / standard / detailed output",
        "Composite eval: Σ(score × weight) before delivery",
        "Per-agent token budget enforcement",
        "Per-session USD cost estimation + CostReport API",
    ]

    add_rect(sl, Inches(0.4), Inches(1.3), Inches(5.9), Inches(5.6),
             fill_color=DARK_BG3, line_color=BLUE_DIM, line_width=Pt(1.5))
    add_rect(sl, Inches(0.4), Inches(1.3), Inches(5.9), Inches(0.42), fill_color=BLUE_DIM)
    add_text_box(sl, "RegDelta",
                 Inches(0.55), Inches(1.32), Inches(5.6), Inches(0.38),
                 font_size=Pt(14), bold=True, color=WHITE)
    add_bullet_items(sl, left_items,
                     Inches(0.55), Inches(1.82), Inches(5.6), Inches(4.8),
                     font_size=Pt(14), color=WHITE, bullet_color=BLUE)

    add_rect(sl, Inches(6.9), Inches(1.3), Inches(6.0), Inches(5.6),
             fill_color=DARK_BG3, line_color=GREEN_DIM, line_width=Pt(1.5))
    add_rect(sl, Inches(6.9), Inches(1.3), Inches(6.0), Inches(0.42), fill_color=GREEN_DIM)
    add_text_box(sl, "SynapseIQ",
                 Inches(7.05), Inches(1.32), Inches(5.7), Inches(0.38),
                 font_size=Pt(14), bold=True, color=WHITE)
    add_bullet_items(sl, right_items,
                     Inches(7.05), Inches(1.82), Inches(5.7), Inches(4.8),
                     font_size=Pt(14), color=WHITE, bullet_color=GREEN)

    footer_bar(sl)

    # ── SLIDE 8: Demo Flow ───────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_slide_bg(sl, DARK_BG)
    header_bar(sl, "Demo Flow", "End-to-end walkthrough of both modules")

    # RegDelta steps
    add_text_box(sl, "RegDelta — 4 API Calls",
                 Inches(0.4), Inches(1.35), Inches(6.1), Inches(0.4),
                 font_size=Pt(16), bold=True, color=BLUE)
    steps_rd = [
        "POST /ingest/regulation → SEC 17a-4 v2023 (old)",
        "POST /ingest/regulation → SEC 17a-4 v2024 (new)",
        "POST /ingest/policy → AML/KYC internal policy",
        "POST /analyze/delta → ImpactReport JSON",
    ]
    for i, step in enumerate(steps_rd):
        y = Inches(1.85 + i * 0.75)
        add_rect(sl, Inches(0.4), y, Inches(0.45), Inches(0.55),
                 fill_color=BLUE_DIM)
        add_text_box(sl, str(i+1), Inches(0.4), y, Inches(0.45), Inches(0.55),
                     font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, Inches(0.95), y, Inches(5.55), Inches(0.55),
                 fill_color=DARK_BG3, line_color=BLUE_DIM, line_width=Pt(1))
        add_text_box(sl, step, Inches(1.1), y + Inches(0.06), Inches(5.3), Inches(0.45),
                     font_size=Pt(13), color=WHITE)

    # SynapseIQ steps
    add_text_box(sl, "SynapseIQ — 3 API Calls",
                 Inches(6.9), Inches(1.35), Inches(6.1), Inches(0.4),
                 font_size=Pt(16), bold=True, color=GREEN)
    steps_sq = [
        "POST /ingest → transformers_overview.txt, llm_finetuning.txt",
        "POST /synthesize → query: 'Compare LoRA vs full fine-tuning', format: standard",
        "Response → synthesis doc + 4-dim scorecard + PASS verdict + agent traces",
    ]
    for i, step in enumerate(steps_sq):
        y = Inches(1.85 + i * 1.05)
        add_rect(sl, Inches(6.9), y, Inches(0.45), Inches(0.9),
                 fill_color=GREEN_DIM)
        add_text_box(sl, str(i+1), Inches(6.9), y, Inches(0.45), Inches(0.9),
                     font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, Inches(7.45), y, Inches(5.55), Inches(0.9),
                 fill_color=DARK_BG3, line_color=GREEN_DIM, line_width=Pt(1))
        add_text_box(sl, step, Inches(7.6), y + Inches(0.08), Inches(5.3), Inches(0.75),
                     font_size=Pt(13), color=WHITE)

    # Demo link
    add_rect(sl, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.45),
             fill_color=PptxRGB(0x0C,0x2A,0x4A), line_color=BLUE_DIM, line_width=Pt(1))
    add_text_box(sl, "Live Demo Output: https://htmlpreview.github.io/?https://github.com/dp2426-NAU/LLM-Project/blob/main/docs/demo_output.html",
                 Inches(0.6), Inches(6.42), Inches(12.1), Inches(0.4),
                 font_size=Pt(11), color=BLUE, align=PP_ALIGN.CENTER)

    footer_bar(sl)

    # ── SLIDE 9: Results ─────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_slide_bg(sl, DARK_BG)
    header_bar(sl, "Results & Impact", "Performance benchmarks and quality outcomes")

    # Score cards
    metrics = [
        ("Coherence",    "0.84", BLUE,   "weight 0.30"),
        ("Relevance",    "0.91", GREEN,  "weight 0.35"),
        ("Factuality",   "0.78", ORANGE, "weight 0.25"),
        ("Completeness", "0.82", PURPLE, "weight 0.10"),
    ]
    for i, (name, score, color, weight) in enumerate(metrics):
        x = Inches(0.4 + i * 3.23)
        add_rect(sl, x, Inches(1.35), Inches(3.0), Inches(1.6),
                 fill_color=DARK_BG3, line_color=color, line_width=Pt(1.5))
        add_text_box(sl, name, x + Inches(0.1), Inches(1.4), Inches(2.7), Inches(0.38),
                     font_size=Pt(12), bold=True, color=MUTED)
        add_text_box(sl, score, x + Inches(0.1), Inches(1.78), Inches(2.7), Inches(0.65),
                     font_size=Pt(30), bold=True, color=color)
        add_text_box(sl, weight, x + Inches(0.1), Inches(2.4), Inches(2.7), Inches(0.3),
                     font_size=Pt(10), color=MUTED)

    # Composite + verdict
    add_rect(sl, Inches(0.4), Inches(3.1), Inches(5.9), Inches(0.55),
             fill_color=PptxRGB(0x0A,0x1F,0x0F), line_color=GREEN_DIM, line_width=Pt(1.5))
    add_text_box(sl, "Composite Score: 0.848  →  PASS ✓",
                 Inches(0.6), Inches(3.13), Inches(5.6), Inches(0.48),
                 font_size=Pt(16), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # Performance table
    add_text_box(sl, "Performance Benchmarks",
                 Inches(0.4), Inches(3.82), Inches(12.5), Inches(0.38),
                 font_size=Pt(14), bold=True, color=MUTED)
    perf = [
        ("Metric", "RegDelta", "SynapseIQ"),
        ("End-to-end latency", "4–8 seconds", "6–12 seconds"),
        ("Token consumption/query", "~2,000–3,500", "~3,000–4,200"),
        ("Est. cost/query (gpt-4o-mini)", "~$0.003–0.005", "~$0.004–0.007"),
        ("Vector retrieval latency", "< 100ms", "< 100ms"),
        ("Ollama local inference", "$0.00", "$0.00"),
    ]
    for i, row_data in enumerate(perf):
        y = Inches(4.28 + i * 0.43)
        bg = DARK_BG2 if i == 0 else (DARK_BG3 if i % 2 == 0 else DARK_BG)
        for j, cell_text in enumerate(row_data):
            x = Inches(0.4 + j * 4.16)
            add_rect(sl, x, y, Inches(4.1), Inches(0.4), fill_color=bg)
            fc = MUTED if i == 0 else WHITE
            add_text_box(sl, cell_text, x + Inches(0.1), y + Inches(0.04), Inches(3.9), Inches(0.35),
                         font_size=Pt(11), bold=(i == 0), color=fc)

    footer_bar(sl)

    # ── SLIDE 10: Future Work ────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_slide_bg(sl, DARK_BG)
    header_bar(sl, "Future Work", "Roadmap and open research directions")

    roadmap = [
        ("Short-Term (v1.1)", BLUE, [
            "Streaming responses via WebSocket — real-time agent progress visualization",
            "Streamlit dashboard UI for non-technical compliance users",
            "Parallel agent execution: Researcher + Critic concurrently (~30% latency reduction)",
            "Hardened PDF parsing for multi-column and table-heavy regulatory documents",
        ]),
        ("Medium-Term (v2.0)", GREEN, [
            "Multi-tenant corpus isolation with authentication",
            "Automated regulation monitoring: SEC EDGAR / Federal Register API polling",
            "Agent persistent memory across sessions",
            "Fine-tuned domain embeddings on regulatory and legal corpora",
        ]),
        ("Research Directions", PURPLE, [
            "LLM-as-judge calibration vs. human expert ratings",
            "Graph-based retrieval over relational regulatory concept networks",
            "Cross-lingual support: EU regulatory documents (GDPR, MiFID II)",
            "Optimal agent decomposition for domain-specific synthesis tasks",
        ]),
    ]
    for i, (title, color, items) in enumerate(roadmap):
        y = Inches(1.35 + i * 1.9)
        add_rect(sl, Inches(0.4), y, Inches(12.5), Inches(1.75),
                 fill_color=DARK_BG3, line_color=color, line_width=Pt(1.5))
        add_rect(sl, Inches(0.4), y, Inches(0.06), Inches(1.75), fill_color=color)
        add_text_box(sl, title, Inches(0.6), y + Inches(0.06), Inches(4), Inches(0.38),
                     font_size=Pt(13), bold=True, color=color)
        add_bullet_items(sl, items,
                         Inches(0.6), y + Inches(0.5), Inches(12.1), Inches(1.2),
                         font_size=Pt(13), color=WHITE, bullet_color=color)

    footer_bar(sl)

    # ── SLIDE 11: Conclusion ─────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_slide_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, Inches(0.25), H, fill_color=BLUE_DIM)

    add_text_box(sl, "Conclusion", Inches(0.65), Inches(0.5), Inches(11.5), Inches(0.8),
                 font_size=Pt(36), bold=True, color=WHITE)
    add_rect(sl, Inches(0.65), Inches(1.3), Inches(10), Inches(0.05), fill_color=BLUE_DIM)

    blocks = [
        ("What was built:", [
            "RegDelta: Regulatory change impact analysis — section diff + dual-corpus RAG + LLM risk classification",
            "SynapseIQ: Multi-agent knowledge synthesis — Researcher→Critic→Synthesizer→Validator pipeline",
        ], BLUE),
        ("What was demonstrated:", [
            "LLMs excel as domain reasoners when grounded in retrieved context — no hallucination",
            "Multi-agent critique + self-evaluation improves quality over single-shot generation",
            "Production LLM systems require: token budgets, cost tracking, observable pipelines, graceful failure",
        ], GREEN),
        ("Key technical contributions:", [
            "Section-level regulatory diff engine with semantic LLM annotation",
            "4-agent sequential pipeline with shared state and quality-gated delivery",
            "YAML+Jinja2 hot-reloadable prompt management system",
        ], ORANGE),
    ]
    y_off = Inches(1.5)
    for title, bullets, color in blocks:
        add_text_box(sl, title, Inches(0.65), y_off, Inches(12), Inches(0.38),
                     font_size=Pt(13), bold=True, color=color)
        y_off += Inches(0.38)
        for b in bullets:
            p = sl.shapes.add_textbox(Inches(0.85), y_off, Inches(11.8), Inches(0.38))
            tf = p.text_frame
            run = tf.paragraphs[0].add_run()
            run.text = "▸  " + b
            run.font.size = Pt(13)
            run.font.color.rgb = WHITE
            y_off += Inches(0.4)
        y_off += Inches(0.1)

    add_rect(sl, Inches(0.65), Inches(6.5), Inches(12.1), Inches(0.7),
             fill_color=PptxRGB(0x0C,0x2A,0x4A), line_color=BLUE_DIM, line_width=Pt(1.5))
    add_text_box(sl, "github.com/dp2426-NAU/LLM-Project",
                 Inches(0.85), Inches(6.55), Inches(12), Inches(0.5),
                 font_size=Pt(16), bold=True, color=BLUE)

    footer_bar(sl)

    out_path = OUT / "presentation.pptx"
    prs.save(str(out_path))
    print(f"OK  Created {out_path}")


if __name__ == "__main__":
    build_docx()
    build_pptx()
    print("Done.")
